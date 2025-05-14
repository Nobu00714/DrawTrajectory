from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import glob
import os
import re

# ベースデータフォルダ
base_folder = '/Volumes/PortableSSD/'

# タスクタイプとパラメータ
task_types = ['CDC', 'CAC', 'Linear', 'Circular', 'Sinewave']
biases = ['Fast', 'Accurate', 'Neutral']
Ws_steering = [16, 22, 36, 90]
Ws_crossing = [8, 15, 28, 50, 100]
A_linear_circular = [360, 480, 700, 850]
A_sinewave = [325, 650, 975, 1300]
A_crossing = [200, 550, 880]

for task_type in task_types:
    for bias in biases:
        if task_type in ['CDC', 'CAC']:
            W_list = Ws_crossing
            A_list = A_crossing
        elif task_type == 'Sinewave':
            W_list = Ws_steering
            A_list = A_sinewave
        else:
            W_list = Ws_steering
            A_list = A_linear_circular

        for W in W_list:
            for A in A_list:
                data_folder = os.path.join(base_folder, f'TrajectoryData_{task_type}')
                pattern = f"*{bias}*A{A}W{W}*.csv"
                csv_files = glob.glob(os.path.join(data_folder, '*', pattern))

                if not csv_files:
                    print(f"スキップ: {task_type}, {bias}, A={A}, W={W} → ファイルなし")
                    continue

                fig, ax = plt.subplots(figsize=(7, 7))

                for file_path in csv_files:
                    df = pd.read_csv(file_path)
                    if 'x座標' in df.columns and 'y座標' in df.columns:
                        ax.plot(df['x座標'], df['y座標'], linewidth=0.2)
                    else:
                        print(f"無視： {file_path} に 'x座標' または 'y座標' がありません。")
                        continue
                if task_type == 'CDC':
                    limit = A / 2 + W + 100
                    ax.set_xlim(-limit, limit)
                    ax.set_ylim(-limit, limit)
                    for x_pos in [-A / 2, A / 2]:
                        rect = patches.Rectangle(
                            (x_pos - 0.5, -W / 2),
                            width=1.0,
                            height=W,
                            color='black',
                            alpha=0.9,
                            zorder=10
                        )
                        ax.add_patch(rect)

                elif task_type == 'CAC':
                    limit = A / 2 + W + 50
                    ax.set_xlim(-limit, limit)
                    ax.set_ylim(-limit, limit)
                    for x_pos in [-A / 2, A / 2]:
                        rect = patches.Rectangle(
                            (x_pos - W / 2, -0.5),
                            width=W,
                            height=1.0,
                            color='black',
                            alpha=0.9,
                            zorder=10
                        )
                        ax.add_patch(rect)

                elif task_type == 'Linear':
                    W_half = W / 2
                    ax.fill_between(
                        x=[-A / 2, A / 2],
                        y1=-W_half,
                        y2=W_half,
                        color='lightgreen',
                        alpha=0.3
                    )
                    ax.set_xlim(-A / 2, A / 2)
                    ax.set_ylim(-A / 2, A / 2)

                elif task_type == 'Circular':
                    r = A / (2 * np.pi)
                    r_inner = r - W / 2
                    r_outer = r + W / 2
                    ax.add_patch(patches.Wedge((0, 0), r_outer, 0, 360, color='lightgreen', alpha=0.5))
                    ax.add_patch(patches.Wedge((0, 0), r_inner, 0, 360, color='white'))
                    ax.set_xlim(-200, 200)
                    ax.set_ylim(-200, 200)

                elif task_type == 'Sinewave':
                    x = np.linspace(-960, 960, 2000)
                    y_center = -40.1 * np.sin(2 * np.pi * 7 * x / 1920)
                    dy_dx = (-40.1 * 2 * np.pi * 7 / 1920) * np.cos(2 * np.pi * 7 * x / 1920)
                    norm = np.sqrt(1 + dy_dx ** 2)
                    nx = -dy_dx / norm
                    ny = 1 / norm

                    x_upper = x + (W / 2) * nx
                    y_upper = y_center + (W / 2) * ny
                    x_lower = x - (W / 2) * nx
                    y_lower = y_center - (W / 2) * ny

                    verts = np.concatenate([
                        np.stack([x_upper, y_upper], axis=1),
                        np.stack([x_lower[::-1], y_lower[::-1]], axis=1)
                    ])
                    ax.add_patch(patches.Polygon(verts, closed=True, color='lightgreen', alpha=0.3))

                    right = 754
                    left = right - (A / 325) * 275
                    ax.set_xlim(left, right)
                    ax.set_ylim(-(right - left) / 2, (right - left) / 2)

                ax.set_autoscale_on(False)
                # ax.set_xlabel('X', fontsize=40)
                # ax.set_ylabel('Y', fontsize=40)
                ax.tick_params(axis='both', labelsize=20)
                ax.set_title(f'{task_type} {bias} A{A} W{W}', fontsize=20)
                ax.grid(False)

                # PNG保存先のフォルダを定義
                png_folder = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'PNG')
                os.makedirs(png_folder, exist_ok=True)  # フォルダがなければ作成

                # 保存パスを指定して保存
                save_filename = os.path.join(png_folder, f"{task_type}{bias}A{A}W{W}.png")
                fig.savefig(save_filename, dpi=300)
                plt.close(fig)

                print(f"保存完了: {save_filename}")

# 1. 読み込みたいPNG画像があるフォルダ
image_folder = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'PNG')

# 2. 出力するPowerPointファイル名
output_pptx = 'All_Trajectories_Visualized.pptx'

# 3. 読み込む画像ファイル一覧を取得
image_files = glob.glob(os.path.join(image_folder, '*.png'))

# 4. ファイル情報をパース
file_info_list = []
for image_path in image_files:
    filename = os.path.basename(image_path)

    # 正規表現でファイル名から情報を抜き出す
    match = re.match(r'(CDC|CAC|Linear|Circular|Sinewave)(Fast|Accurate|Neutral)A(\d+)W(\d+)\.png', filename)
    if match:
        task, bias, A, W = match.groups()
        A = int(A)
        W = int(W)
        file_info_list.append((task, A, W, bias, image_path))
    else:
        print(f"スキップ：ファイル名形式が違うため無視: {filename}")

# 5. ソート
# Steering順: Linear → Circular → Sinewave
steering_order = {'CDC': 0, 'CAC': 1, 'Linear': 2, 'Circular': 3, 'Sinewave': 4}
# Bias順: Accurate → Neutral → Fast に修正
bias_order = {'Accurate': 0, 'Neutral': 1, 'Fast': 2}

sorted_files = sorted(
    file_info_list,
    key=lambda x: (steering_order[x[0]], x[1], x[2], bias_order[x[3]])
)

# 6. 新しいプレゼンテーションを作成
prs = Presentation()

# 7. 並べた順にスライド追加
for task, A, W, bias, image_path in sorted_files:
    slide_layout = prs.slide_layouts[6]  # 白紙スライド
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.5)
    top = Inches(0.5)
    height = Inches(7)
    slide.shapes.add_picture(image_path, left, top, height=height)

# 8. 保存
prs.save(output_pptx)

print(f"パワポ作成完了: {output_pptx}")
